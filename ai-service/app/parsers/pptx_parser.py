"""PowerPoint Parser using python-pptx to extract slide text and embedded images slide by slide."""

import io
import base64
from typing import List
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from app.models import ExtractedPage, ExtractionResult


def parse_pptx(file_path_or_bytes: bytes | str) -> ExtractionResult:
    """Extract text and images from a PPTX file per slide.
    
    Args:
        file_path_or_bytes: File path or raw bytes of the PPTX file.
        
    Returns:
        ExtractionResult containing slide-by-slide text and base64 images.
    """
    if isinstance(file_path_or_bytes, bytes):
        prs = Presentation(io.BytesIO(file_path_or_bytes))
    else:
        prs = Presentation(file_path_or_bytes)

    extracted_pages: List[ExtractedPage] = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_text_parts = []
        slide_images = []

        # Extract text from slide shapes and notes
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    p_text = paragraph.text.strip()
                    if p_text:
                        slide_text_parts.append(p_text)

            # Extract embedded picture shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                ext = image.ext
                b64_str = base64.b64encode(image_bytes).decode("utf-8")
                data_uri = f"data:image/{ext};base64,{b64_str}"
                slide_images.append(data_uri)

        # Also extract speaker notes if available
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_text_parts.append(f"[Speaker Notes: {notes}]")

        combined_text = "\n".join(slide_text_parts)

        extracted_pages.append(
            ExtractedPage(
                page_number=slide_idx + 1,
                text=combined_text,
                images=slide_images
            )
        )

    return ExtractionResult(
        file_type="pptx",
        total_pages=len(extracted_pages),
        pages=extracted_pages,
        metadata={"total_slides": len(extracted_pages)}
    )
